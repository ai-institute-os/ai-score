"""
Generate the AI Institute 13-slide strategic presentation as a .pptx file.
Returns the file as bytes for streaming via FastAPI.

Usage:
    from scripts.generate_ai_institute_pptx import build_pptx
    pptx_bytes = build_pptx()
"""
from __future__ import annotations

import io
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x12, 0x67)
NAVY_MID   = RGBColor(0x0F, 0x14, 0x70)
ACCENT     = RGBColor(0xFF, 0x6B, 0x35)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT       = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_MID   = RGBColor(0x33, 0x41, 0x55)
TEXT_LIGHT = RGBColor(0x64, 0x74, 0x8B)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
BG_LIGHT   = RGBColor(0xF8, 0xF9, 0xFC)
DARK_GREEN = RGBColor(0x16, 0x65, 0x34)  # dark forest green for accent borders

# ── Slide dimensions (16:9 widescreen) ─────────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)

# ── Margins ────────────────────────────────────────────────────────────────────
LM   = Inches(1.1)    # left margin
RM   = Inches(1.1)    # right margin
TM   = Inches(0.75)   # top margin for label
CW   = SW - LM - RM  # content width


# ── Low-level helpers ──────────────────────────────────────────────────────────

def _blank(prs: Presentation) -> object:
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def _bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_line(shape) -> None:
    """Remove the border/outline from a shape via direct XML manipulation."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    for existing in spPr.findall(qn('a:ln')):
        spPr.remove(existing)
    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', '0')
    etree.SubElement(ln, qn('a:noFill'))


def _rect(slide, left, top, width, height, fill_color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    _no_line(shape)
    f = shape.fill
    f.solid()
    f.fore_color.rgb = fill_color
    return shape


def _txt(slide, text: str, left, top, width, height, *,
         size: float = 18,
         bold: bool = False,
         italic: bool = False,
         color: RGBColor = TEXT,
         align: PP_ALIGN = PP_ALIGN.LEFT,
         font: str = "Inter",
         wrap: bool = True,
         line_spacing: float | None = None) -> object:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as Pt_
        from pptx.oxml.ns import qn
        from lxml import etree
        # line spacing via XML
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))
    return txb


def _para(txb, text: str, *,
          size: float = 14,
          bold: bool = False,
          color: RGBColor = TEXT_MID,
          align: PP_ALIGN = PP_ALIGN.LEFT,
          font: str = "Inter",
          space_before: float = 0) -> None:
    tf = txb.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _label(slide, text: str, dark: bool = True) -> None:
    color = ACCENT if dark else ACCENT
    _txt(slide, text, LM, TM, CW, Inches(0.4),
         size=9, bold=True, color=color, font="Inter")


def _accent_rule(slide, top, width=Inches(0.5), dark: bool = True) -> None:
    _rect(slide, LM, top, width, Pt(3), ACCENT)


def _title(slide, text: str, top, *, dark: bool = True, size: float = 36) -> object:
    color = WHITE if dark else NAVY
    return _txt(slide, text, LM, top, CW, Inches(1.6),
                size=size, bold=True, color=color, font="Playfair Display" if False else "Inter",
                wrap=True)


# ── SLIDE BUILDERS ─────────────────────────────────────────────────────────────

def _s01_cover(prs):
    slide = _blank(prs)
    _bg(slide, NAVY)

    # Top brand
    _txt(slide, "AIScore™  by AI Institute ApS",
         LM, Inches(0.5), CW, Inches(0.5),
         size=13, color=RGBColor(0xCC, 0xCC, 0xDD), font="Inter")

    # Eyebrow
    _txt(slide, "STRATEGISK MØDE — FORTROLIGT",
         LM, Inches(1.6), CW, Inches(0.4),
         size=9, bold=True, color=ACCENT, font="Inter")

    # Big title
    tb = _txt(slide, "AI Institute",
              LM, Inches(2.1), CW, Inches(2.0),
              size=72, bold=True, color=WHITE, font="Inter")

    # Subtitle
    _txt(slide, "Hvordan AI-systemer opfatter, kategoriserer og vælger virksomheder",
         LM, Inches(4.2), Inches(8), Inches(1.0),
         size=18, color=RGBColor(0xCC, 0xCC, 0xEE), font="Inter")

    # Product chips row
    chips = ["AIScore", "InsideAI", "AISelect", "RunAI"]
    x = LM
    y = Inches(5.9)
    _rect(slide, LM - Inches(0.1), Inches(5.75), CW + Inches(0.2), Inches(0.02),
          RGBColor(0x33, 0x38, 0x88))
    for chip in chips:
        _txt(slide, chip, x, y + Inches(0.1), Inches(1.5), Inches(0.35),
             size=11, color=RGBColor(0xAA, 0xAA, 0xCC), font="Inter")
        x += Inches(1.65)

    return slide


def _s02_purpose(prs):
    slide = _blank(prs)
    _bg(slide, WHITE)
    _label(slide, "02 — FORMÅL", dark=False)

    _txt(slide, "Hvorfor viser jeg dig det her?",
         LM, Inches(1.3), CW, Inches(1.0),
         size=36, bold=True, color=NAVY, font="Inter")

    lines = [
        "Jeg leder ikke efter feedback på idéen.",
        "Jeg leder ikke efter en investor.",
        "Jeg leder ikke efter en kunde.",
        "",
        "Jeg forsøger at finde den rette tekniske arkitektur",
        "og de rette mennesker til at bygge dette korrekt.",
        "",
        "Det kan være dig.",
        "Det kan være nogen i dit netværk.",
    ]
    y = Inches(2.8)
    for i, line in enumerate(lines):
        if line == "":
            y += Inches(0.15)
            continue
        bold = i in (4, 5)
        color = NAVY if bold else TEXT_MID
        size = 15 if bold else 14
        _txt(slide, line, LM, y, Inches(9), Inches(0.4),
             size=size, bold=bold, color=color, font="Inter")
        y += Inches(0.42)

    return slide


def _s03_shift(prs):
    slide = _blank(prs)
    _bg(slide, BG_LIGHT)
    _label(slide, "03 — KONTEKST", dark=False)

    _txt(slide, "Det skifte de fleste overser",
         LM, Inches(1.3), CW, Inches(0.9),
         size=34, bold=True, color=NAVY, font="Inter")

    # Left card (past / passive) — navy left border + thin bottom bar
    card_w = Inches(5.5)
    _rect(slide, LM, Inches(2.5), card_w, Inches(4.0), WHITE)
    _rect(slide, LM, Inches(2.5), Inches(0.06), Inches(4.0), NAVY)
    _rect(slide, LM, Inches(6.47), card_w, Inches(0.03), NAVY)
    _txt(slide, "FORTID", LM + Inches(0.2), Inches(2.6), card_w, Inches(0.35),
         size=9, bold=True, color=TEXT_LIGHT, font="Inter")
    _txt(slide, "Google rangerede virksomheder.",
         LM + Inches(0.2), Inches(3.0), card_w - Inches(0.4), Inches(0.85),
         size=22, bold=True, color=NAVY, font="Inter")
    _txt(slide, "Virksomheder konkurrerede om placeringer.\nSynlighed var en rangliste.",
         LM + Inches(0.2), Inches(4.0), card_w - Inches(0.4), Inches(0.9),
         size=13, color=TEXT, font="Inter")

    # Right card (now / active) — orange left border + thin bottom bar
    rx = LM + card_w + Inches(0.2)
    _rect(slide, rx, Inches(2.5), card_w, Inches(4.0), WHITE)
    _rect(slide, rx, Inches(6.47), card_w, Inches(0.03), ACCENT)
    _rect(slide, rx, Inches(2.5), Inches(0.06), Inches(4.0), ACCENT)  # orange left border
    _txt(slide, "NU", rx + Inches(0.25), Inches(2.6), card_w, Inches(0.35),
         size=9, bold=True, color=ACCENT, font="Inter")
    _txt(slide, "AI-systemer vælger virksomheder.",
         rx + Inches(0.25), Inches(3.0), card_w - Inches(0.5), Inches(0.85),
         size=22, bold=True, color=NAVY, font="Inter")
    _txt(slide, "Kunden får et svar.\nOg i det svar er virksomheden enten valgt — eller fravalgt.",
         rx + Inches(0.25), Inches(4.0), card_w - Inches(0.5), Inches(1.0),
         size=13, color=TEXT_MID, font="Inter")

    return slide


def _s04_sentence(prs):
    slide = _blank(prs)
    _bg(slide, WHITE)
    _label(slide, "04 — KERNEN", dark=False)

    _txt(slide, "AI Institute på én sætning",
         LM, Inches(1.3), CW, Inches(0.9),
         size=32, bold=True, color=NAVY, font="Inter")

    _txt(slide, "De fleste virksomheder forsøger at blive fundet.",
         LM, Inches(2.6), CW, Inches(0.7),
         size=20, color=TEXT_LIGHT, font="Inter")

    _txt(slide, "AI Institute forsøger at forstå,\nhvorfor virksomheder bliver valgt.",
         LM, Inches(3.3), CW, Inches(1.3),
         size=28, bold=True, color=NAVY, font="Inter")

    # Product grid — no outlines; last cell (RunAI) gets dark green left border
    products = [
        ("AIScore",  "MÅLER",     "Kortlægger AI-positioner"),
        ("InsideAI", "OVERVÅGER", "Tracker ændringer over tid"),
        ("AISelect", "PÅVIRKER",  "Arbejder med positionering"),
        ("RunAI",    "UDFØRER",   "Implementerer i organisationen"),
    ]
    cell_w = CW / 4
    gx = LM
    gy = Inches(5.0)
    for i, (name, verb, desc) in enumerate(products):
        is_last = (i == 3)
        _rect(slide, gx, gy, cell_w - Inches(0.05), Inches(1.9), BG_LIGHT)
        if is_last:
            _rect(slide, gx, gy, Inches(0.06), Inches(1.9), DARK_GREEN)
        else:
            _rect(slide, gx, gy, Inches(0.06), Inches(1.9), NAVY)
        _txt(slide, verb, gx + Inches(0.2), gy + Inches(0.15), cell_w - Inches(0.3), Inches(0.35),
             size=9, bold=True, color=ACCENT, font="Inter")
        _txt(slide, name, gx + Inches(0.2), gy + Inches(0.5), cell_w - Inches(0.3), Inches(0.4),
             size=14, bold=True, color=NAVY, font="Inter")
        _txt(slide, desc, gx + Inches(0.2), gy + Inches(0.9), cell_w - Inches(0.3), Inches(0.8),
             size=11, color=TEXT_LIGHT, font="Inter")
        gx += cell_w

    return slide


def _s05_comparison(prs):
    slide = _blank(prs)
    _bg(slide, BG_LIGHT)
    _label(slide, "05 — DIFFERENTIERING", dark=False)

    _txt(slide, "Synlighed er ikke slutmålet. Udvælgelse er.",
         LM, Inches(1.3), CW, Inches(0.85),
         size=30, bold=True, color=NAVY, font="Inter")

    # Left col: GEO
    col_w = Inches(5.3)
    cy = Inches(2.4)
    _rect(slide, LM, cy, col_w, Inches(4.0), WHITE)
    _rect(slide, LM, cy, Inches(0.06), Inches(4.0), DARK_GREEN)
    _txt(slide, "GEO / AI SEO", LM + Inches(0.2), cy + Inches(0.15), col_w, Inches(0.4),
         size=14, bold=True, color=TEXT_LIGHT, font="Inter")
    _txt(slide, "Fokus: At blive nævnt", LM + Inches(0.2), cy + Inches(0.65), col_w, Inches(0.35),
         size=12, color=TEXT_LIGHT, font="Inter")
    qs_l = ["Bliver vi nævnt?", "Bliver vi citeret?", "Er vi synlige?"]
    qy = cy + Inches(1.15)
    for q in qs_l:
        _txt(slide, "→  " + q, LM + Inches(0.2), qy, col_w, Inches(0.35),
             size=12, color=TEXT_LIGHT, font="Inter")
        qy += Inches(0.38)
    _txt(slide, "Mål: Mere omtale", LM + Inches(0.2), cy + Inches(3.35), col_w, Inches(0.4),
         size=12, bold=True, color=TEXT_LIGHT, font="Inter")

    # Right col: AI Institute
    rx = LM + col_w + Inches(0.3)
    _rect(slide, rx, cy, col_w, Inches(4.0), WHITE)
    _rect(slide, rx, cy, col_w, Inches(0.06), NAVY)
    _rect(slide, rx, cy, Inches(0.06), Inches(4.0), ACCENT)
    _txt(slide, "AI Institute", rx + Inches(0.25), cy + Inches(0.15), col_w, Inches(0.4),
         size=14, bold=True, color=NAVY, font="Inter")
    _txt(slide, "Fokus: At blive valgt", rx + Inches(0.25), cy + Inches(0.65), col_w, Inches(0.35),
         size=12, color=TEXT_MID, font="Inter")
    qs_r = ["Hvorfor vælger AI én virksomhed?", "Hvilken rolle har AI tildelt virksomheden?",
            "Hvad driver AI-udvælgelse?", "Hvornår fravælges vi?"]
    qy = cy + Inches(1.15)
    for q in qs_r:
        _txt(slide, "→  " + q, rx + Inches(0.25), qy, col_w - Inches(0.4), Inches(0.38),
             size=12, color=TEXT_MID, font="Inter")
        qy += Inches(0.42)
    _txt(slide, "Mål: Forstå og påvirke AI-valg", rx + Inches(0.25), cy + Inches(3.35), col_w, Inches(0.4),
         size=12, bold=True, color=NAVY, font="Inter")

    return slide


def _s06_problem(prs):
    slide = _blank(prs)
    _bg(slide, NAVY)
    _label(slide, "06 — PROBLEM", dark=True)

    _txt(slide, "Hvorfor det betyder noget",
         LM, Inches(1.3), CW, Inches(0.85),
         size=34, bold=True, color=WHITE, font="Inter")

    _txt(slide, "De fleste virksomheder ved ikke:",
         LM, Inches(2.6), CW, Inches(0.45),
         size=15, color=RGBColor(0xCC, 0xCC, 0xEE), font="Inter")

    items = [
        "hvordan AI ser dem",
        "hvilken kategorirolle AI har givet dem",
        "hvornår de vælges",
        "hvornår de fravælges",
        "hvorfor konkurrenter vælges i stedet",
        "hvordan deres AI-position ændrer sig over tid",
    ]
    y = Inches(3.1)
    for item in items:
        _rect(slide, LM, y + Inches(0.06), Inches(0.06), Inches(0.28), ACCENT)
        _txt(slide, item, LM + Inches(0.25), y, Inches(9), Inches(0.45),
             size=15, color=RGBColor(0xCC, 0xCC, 0xEE), font="Inter")
        y += Inches(0.55)

    return slide


def _s07_research(prs):
    slide = _blank(prs)
    _bg(slide, NAVY_MID)
    _label(slide, "07 — FUNDAMENTET", dark=True)

    _txt(slide, "AI Institute",
         LM, Inches(1.5), CW, Inches(1.3),
         size=56, bold=True, color=WHITE, font="Inter")
    _txt(slide, "FORSKNINGSMOTOREN",
         LM, Inches(3.0), CW, Inches(0.5),
         size=12, bold=True, color=ACCENT, font="Inter")

    _txt(slide, "AI Institute eksisterer for at forstå, hvordan AI-systemer\nopfatter, kategoriserer og vælger virksomheder.",
         LM, Inches(4.0), Inches(9), Inches(1.0),
         size=16, color=RGBColor(0xCC, 0xCC, 0xEE), font="Inter")

    _txt(slide, "Produkterne er ikke separate idéer.\nDe er konsekvenser af den samme forskning.",
         LM, Inches(5.3), Inches(9), Inches(0.9),
         size=15, color=RGBColor(0xAA, 0xAA, 0xCC), font="Inter")

    return slide


def _s08_timing(prs):
    slide = _blank(prs)
    _bg(slide, NAVY)
    _label(slide, "08 — TIMING", dark=True)

    _txt(slide, "Hvorfor timing betyder noget",
         LM, Inches(1.3), CW, Inches(0.85),
         size=34, bold=True, color=WHITE, font="Inter")

    years = [
        ("2024", "BEGYNDELSEN", "AI begynder at anbefale virksomheder. Positioner er plastiske."),
        ("2025", "OPDAGELSE",   "Virksomheder opdager problemet. Tidligt vindue for at handle."),
        ("2026", "NU",          "AI-positioner formes aktivt. Den mest kritiske fase."),
        ("2027+", "KONSOLIDERING", "Kategoriopfattelser sætter sig. Vanskeligere at ændre."),
    ]
    col_w = CW / 4
    cx = LM
    cy = Inches(2.4)
    for i, (yr, status, body) in enumerate(years):
        is_now = (i == 2)
        bg_col = RGBColor(0x12, 0x18, 0x75) if not is_now else RGBColor(0x18, 0x1F, 0x85)
        _rect(slide, cx, cy, col_w - Inches(0.08), Inches(3.5), bg_col)
        if is_now:
            _rect(slide, cx, cy, Inches(0.06), Inches(3.5), ACCENT)
        yr_color = ACCENT if is_now else WHITE
        _txt(slide, yr, cx + Inches(0.2), cy + Inches(0.2), col_w, Inches(0.6),
             size=28, bold=True, color=yr_color, font="Inter")
        _txt(slide, status, cx + Inches(0.2), cy + Inches(0.85), col_w, Inches(0.35),
             size=9, bold=True, color=ACCENT, font="Inter")
        _txt(slide, body, cx + Inches(0.2), cy + Inches(1.3), col_w - Inches(0.35), Inches(1.9),
             size=12, color=RGBColor(0xAA, 0xAA, 0xCC), font="Inter")
        cx += col_w

    _rect(slide, LM, Inches(6.1), CW, Inches(0.04), RGBColor(0xFF, 0x6B, 0x35))
    _txt(slide, "De virksomheder der etablerer stærke AI-positioner tidligt, får en strukturel fordel senere.",
         LM, Inches(6.25), CW, Inches(0.7),
         size=13, bold=True, color=RGBColor(0xCC, 0xCC, 0xEE), font="Inter")

    return slide


def _s09_aiscore(prs):
    slide = _blank(prs)
    _bg(slide, WHITE)
    _label(slide, "09 — BEVIS", dark=False)

    _txt(slide, "AIScore — Diagnosen",
         LM, Inches(1.3), CW, Inches(0.85),
         size=34, bold=True, color=NAVY, font="Inter")

    _txt(slide, "AIScore viser, hvordan AI-systemer aktuelt positionerer en virksomhed.",
         LM, Inches(2.6), Inches(6.5), Inches(0.55),
         size=14, color=TEXT_MID, font="Inter")

    _txt(slide, "Rapporten svarer på:",
         LM, Inches(3.25), Inches(6.5), Inches(0.35),
         size=13, color=TEXT_MID, font="Inter")

    bullets = [
        "om virksomheden vælges",
        "hvornår den vælges / fravælges",
        "hvilken rolle AI har givet virksomheden",
        "hvilke strukturelle mønstre der begrænser udvælgelse",
    ]
    by = Inches(3.7)
    for b in bullets:
        _rect(slide, LM, by + Inches(0.1), Inches(0.06), Inches(0.22), ACCENT)
        _txt(slide, b, LM + Inches(0.22), by, Inches(6.2), Inches(0.4),
             size=13, color=TEXT_MID, font="Inter")
        by += Inches(0.48)

    # Report preview box
    rx = LM + Inches(6.8)
    rw = Inches(4.5)
    _rect(slide, rx, Inches(2.3), rw, Inches(0.55), NAVY)
    _txt(slide, "AIScore Rapport — Nailster A/S",
         rx + Inches(0.2), Inches(2.38), rw - Inches(0.4), Inches(0.35),
         size=11, color=WHITE, font="Inter")

    pages = [
        ("Side 1", "Forside"),
        ("Side 5", "Executive Summary"),
        ("Side 9", "Strukturelle svagheder"),
        ("Side 12", "Strategisk position"),
    ]
    pw = rw / 2 - Inches(0.15)
    ph = Inches(1.9)
    px, py = rx + Inches(0.1), Inches(2.95)
    for i, (pg, name) in enumerate(pages):
        col = i % 2
        row = i // 2
        bx = px + col * (pw + Inches(0.1))
        by = py + row * (ph + Inches(0.1))
        _rect(slide, bx, by, pw, ph, BG_LIGHT)
        _txt(slide, pg,   bx + Inches(0.1), by + Inches(0.15), pw, Inches(0.3),
             size=9, color=TEXT_LIGHT, font="Inter")
        _txt(slide, name, bx + Inches(0.1), by + Inches(0.5), pw, Inches(0.65),
             size=11, bold=True, color=NAVY, font="Inter")
        _rect(slide, bx + Inches(0.1), by + Inches(1.55), Inches(0.4), Inches(0.04), ACCENT)

    return slide


def _s10_flow(prs):
    slide = _blank(prs)
    _bg(slide, BG_LIGHT)
    _label(slide, "10 — SYSTEMET", dark=False)

    _txt(slide, "Fra diagnose til implementering",
         LM, Inches(1.3), CW, Inches(0.85),
         size=34, bold=True, color=NAVY, font="Inter")

    steps = [
        ("AIScore",  "DIAGNOSE",    "Kortlægger AI-positioner og identificerer strukturelle svagheder"),
        ("InsideAI", "OVERVÅGNING", "Tracker løbende ændringer og observerer udviklingen"),
        ("AISelect", "LØSNING",     "Arbejder med AI-positionering og påvirker AI-valg"),
    ]
    sy = Inches(2.65)
    for i, (name, tag, desc) in enumerate(steps):
        border_col = ACCENT if i == 0 else BORDER
        _rect(slide, LM, sy, Inches(5.5), Inches(1.15), WHITE)
        _rect(slide, LM, sy, Inches(0.06), Inches(1.15), border_col)
        _txt(slide, tag, LM + Inches(0.2), sy + Inches(0.1), Inches(5), Inches(0.32),
             size=9, bold=True, color=ACCENT, font="Inter")
        _txt(slide, name, LM + Inches(0.2), sy + Inches(0.42), Inches(5), Inches(0.38),
             size=15, bold=True, color=NAVY, font="Inter")
        _txt(slide, desc, LM + Inches(0.2), sy + Inches(0.78), Inches(5), Inches(0.35),
             size=11, color=TEXT_LIGHT, font="Inter")
        if i < 2:
            _txt(slide, "↓", LM + Inches(2.5), sy + Inches(1.2), Inches(0.5), Inches(0.45),
                 size=16, color=TEXT_LIGHT, font="Inter", align=PP_ALIGN.CENTER)
        sy += Inches(1.6)

    # Right text
    rx = LM + Inches(6.2)
    rlines = [
        ("AIScore identificerer problemet.", False),
        ("InsideAI overvåger udviklingen.", False),
        ("AISelect arbejder med løsningen.", False),
        ("Produkterne er ikke siloer.\nDe er faser i det samme system.", True),
    ]
    ry = Inches(2.65)
    for text, highlight in rlines:
        if highlight:
            _rect(slide, rx, ry + Inches(0.3), Inches(5.5), Inches(0.85), WHITE)
            _rect(slide, rx, ry + Inches(0.3), Inches(0.06), Inches(0.85), ACCENT)
            _txt(slide, text, rx + Inches(0.2), ry + Inches(0.38), Inches(5.2), Inches(0.7),
                 size=13, bold=True, color=NAVY, font="Inter")
        else:
            _txt(slide, text, rx, ry, Inches(5.5), Inches(0.4),
                 size=14, color=TEXT_MID, font="Inter")
        ry += Inches(0.65) if not highlight else Inches(1.2)

    return slide


def _s11_ecosystem(prs):
    slide = _blank(prs)
    _bg(slide, NAVY)
    _label(slide, "11 — STRUKTUREN", dark=True)

    _txt(slide, "Økosystemet",
         LM, Inches(1.3), CW, Inches(0.85),
         size=34, bold=True, color=WHITE, font="Inter")

    items = [
        ("AI Institute", "KERNEN"),
        ("AIScore",      "MÅLER"),
        ("InsideAI",     "OVERVÅGER"),
        ("AISelect",     "PÅVIRKER"),
        ("RunAI",        "UDFØRER"),
    ]
    cx = SW / 2
    ew = Inches(4.0)
    ey = Inches(2.2)
    for i, (name, tag) in enumerate(items):
        is_hub = (i == 0)
        fill = WHITE if is_hub else RGBColor(0x14, 0x1A, 0x78)
        border = WHITE if is_hub else RGBColor(0x33, 0x38, 0x88)
        _rect(slide, cx - ew / 2, ey, ew, Inches(0.65), fill)
        nc = NAVY if is_hub else WHITE
        tc = NAVY if is_hub else ACCENT
        _txt(slide, name, cx - ew / 2 + Inches(0.2), ey + Inches(0.05), ew, Inches(0.35),
             size=15, bold=True, color=nc, font="Inter", align=PP_ALIGN.CENTER)
        _txt(slide, tag, cx - ew / 2 + Inches(0.2), ey + Inches(0.38), ew, Inches(0.22),
             size=9, bold=True, color=tc, font="Inter", align=PP_ALIGN.CENTER)
        ey += Inches(0.65)
        if i < 4:
            _txt(slide, "↓", cx - Inches(0.15), ey + Inches(0.0), Inches(0.3), Inches(0.35),
                 size=14, color=RGBColor(0x44, 0x49, 0x99), font="Inter", align=PP_ALIGN.CENTER)
            ey += Inches(0.35)

    return slide


def _s12_moat(prs):
    slide = _blank(prs)
    _bg(slide, NAVY_MID)
    _label(slide, "12 — MOAT", dark=True)

    _txt(slide, "Hvorfor økosystemet\nbliver stærkere over tid",
         LM, Inches(1.3), Inches(6), Inches(1.4),
         size=34, bold=True, color=WHITE, font="Inter")

    steps = [
        "AI Institute forsker",
        "AIScore producerer rapporter",
        "InsideAI akkumulerer data",
        "AISelect leverer resultater",
        "Mere data → bedre forståelse",
        "Bedre forståelse → bedre produkter",
    ]
    sy = Inches(3.0)
    for s in steps:
        _rect(slide, LM, sy + Inches(0.12), Inches(0.1), Inches(0.1), ACCENT)
        _txt(slide, s, LM + Inches(0.3), sy, Inches(5.5), Inches(0.4),
             size=14, bold=True, color=WHITE, font="Inter")
        sy += Inches(0.52)

    # Moat box
    _rect(slide, LM + Inches(6.5), Inches(2.5), Inches(5.5), Inches(2.8), RGBColor(0x0D, 0x10, 0x55))
    _rect(slide, LM + Inches(6.5), Inches(2.5), Inches(0.08), Inches(2.8), ACCENT)
    _txt(slide, "MOAT", LM + Inches(6.7), Inches(2.65), Inches(5.0), Inches(0.35),
         size=9, bold=True, color=ACCENT, font="Inter")
    _txt(slide, "Moat'en er ikke én rapport.\nMoat'en er den akkumulerede\nforståelse af, hvorfor AI vælger.",
         LM + Inches(6.7), Inches(3.1), Inches(5.0), Inches(1.8),
         size=15, bold=True, color=WHITE, font="Inter")

    return slide


def _s13_ask(prs):
    slide = _blank(prs)
    _bg(slide, NAVY)
    _label(slide, "13 — ANMODNINGEN", dark=True)

    _txt(slide, "Hvor jeg sidder fast",
         LM, Inches(1.3), CW, Inches(0.85),
         size=34, bold=True, color=WHITE, font="Inter")

    negatives = [
        "Jeg mangler ikke idéer.",
        "Jeg mangler ikke produkter.",
        "Jeg mangler ikke salg.",
    ]
    ny = Inches(2.6)
    for neg in negatives:
        _txt(slide, neg, LM, ny, Inches(5.5), Inches(0.4),
             size=14, color=RGBColor(0x88, 0x88, 0xAA), font="Inter")
        ny += Inches(0.45)

    _txt(slide, "Jeg mangler den tekniske arkitektur og de rette\nmennersker til at bygge systemet korrekt.",
         LM, Inches(4.05), Inches(6.0), Inches(1.0),
         size=16, bold=True, color=WHITE, font="Inter")

    needs = [
        "Systemarkitektur",
        "Agent orchestration",
        "Data infrastructure",
        "Teknisk lederskab",
        "QA og governance",
        "Langsigtet technical ownership",
    ]
    nw = Inches(2.5)
    nx, nny = LM + Inches(6.8), Inches(2.5)
    for i, need in enumerate(needs):
        col = i % 2
        row = i // 2
        bx = nx + col * (nw + Inches(0.15))
        by = nny + row * Inches(0.72)
        _rect(slide, bx, by, nw, Inches(0.58), RGBColor(0x14, 0x1A, 0x78))
        _txt(slide, "→  " + need, bx + Inches(0.1), by + Inches(0.1), nw - Inches(0.15), Inches(0.4),
             size=11, color=RGBColor(0xCC, 0xCC, 0xEE), font="Inter")

    _rect(slide, LM, Inches(5.35), CW, Inches(0.03), RGBColor(0x33, 0x38, 0x88))
    _txt(slide, "Det er derfor jeg gerne vil tale med dig.\nMåske er du den rigtige person. Måske kender du den rigtige person.",
         LM, Inches(5.55), Inches(9), Inches(1.0),
         size=15, bold=True, color=RGBColor(0xCC, 0xCC, 0xEE), font="Inter")

    return slide


# ── Public build function ──────────────────────────────────────────────────────

def build_pptx() -> bytes:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    builders = [
        _s01_cover, _s02_purpose, _s03_shift,  _s04_sentence,
        _s05_comparison, _s06_problem, _s07_research, _s08_timing,
        _s09_aiscore, _s10_flow, _s11_ecosystem, _s12_moat, _s13_ask,
    ]
    for fn in builders:
        fn(prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


if __name__ == "__main__":
    data = build_pptx()
    with open("ai-institute-deck.pptx", "wb") as f:
        f.write(data)
    print(f"Saved {len(data):,} bytes → ai-institute-deck.pptx")
